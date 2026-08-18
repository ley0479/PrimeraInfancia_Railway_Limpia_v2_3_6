"""
Servicios puros del Calendario Inteligente de Entregables.

Alpha15: reglas de estado, color, normalización de columnas y conversión de fechas
sin depender de Flask para facilitar pruebas y evitar acoplamiento con otros módulos.
"""
from __future__ import annotations

import calendar
import hashlib
import re
import unicodedata
from datetime import date, datetime, timedelta
from typing import Any, Iterable

try:
    import pandas as pd
except Exception:  # pragma: no cover
    pd = None

ESTADOS_PERMITIDOS = [
    "programado",
    "pendiente",
    "proximo_a_vencer",
    "vence_hoy",
    "vencido",
    "entregado",
    "aprobado",
    "rechazado",
    "no_aplica",
    "cerrado",
]

MODULOS_PERMITIDOS = [
    "RPP",
    "Bienestarina",
    "RAM/RAN/Asistencia",
    "Nutrición",
    "Talento Humano",
    "Planeación Pedagógica",
    "Gestión por Coordinador",
    "Reportes Gerenciales",
    "Cumplimiento ICBF",
]

RECURRENCIAS_PERMITIDAS = ["ninguna", "diaria", "semanal", "mensual"]

COLOR_BY_STATUS = {
    "programado": "azul",
    "pendiente": "azul",
    "proximo_a_vencer": "amarillo",
    "vence_hoy": "rojo",
    "vencido": "rojo",
    "entregado": "verde",
    "aprobado": "verde",
    "rechazado": "rojo",
    "no_aplica": "gris",
    "cerrado": "gris",
}

COLUMN_SYNONYMS = {
    "fecha_limite": ["fecha", "fecha limite", "fecha límite", "fecha de entrega", "fecha entrega", "vence", "vencimiento"],
    "fecha_inicio": ["fecha inicio", "inicio", "desde"],
    "titulo": ["actividad", "entregable", "tarea", "titulo", "título", "nombre actividad", "nombre del entregable"],
    "descripcion": ["descripcion", "descripción", "detalle", "observacion", "observación", "observaciones"],
    "modulo": ["modulo", "módulo", "area", "área", "proceso", "menu", "menú"],
    "tipo_formato": ["formato", "tipo formato", "tipo de formato", "documento", "plantilla"],
    "responsable_nombre": ["responsable", "encargado", "responsable nombre", "responsable_nombre", "persona responsable", "th a cargo", "talento humano a cargo"],
    "coordinador": ["coordinador", "coordinador responsable", "coordinadora", "coord"],
    "unidad": ["unidad", "uds", "sede", "unidad de servicio", "unidad servicio", "nombre uds", "comunidad"],
    "municipio": ["municipio", "ciudad", "localidad"],
    "prioridad": ["prioridad", "nivel", "urgencia"],
    "estado": ["estado", "estatus", "situacion", "situación"],
    "observaciones": ["observacion", "observación", "observaciones", "nota", "notas", "comentario"],
    "componente": ["componente", "area componente", "componente institucional"],
    "entregables": ["entrega", "entregas", "entregable requerido", "evidencia requerida", "soportes"],
    "numero": ["n", "n°", "numero", "número", "item", "ítem"],
}


def normalizar_texto(value: Any) -> str:
    texto = str(value or "").strip().lower()
    texto = unicodedata.normalize("NFKD", texto)
    texto = "".join(ch for ch in texto if not unicodedata.combining(ch))
    texto = texto.replace("ñ", "n")
    texto = re.sub(r"[^a-z0-9]+", " ", texto)
    return " ".join(texto.split())


def parse_fecha(value: Any) -> str | None:
    """Convierte fechas de Excel/CSV a ISO YYYY-MM-DD."""
    if value is None or value == "":
        return None
    if isinstance(value, datetime):
        return value.date().isoformat()
    if isinstance(value, date):
        return value.isoformat()
    text = str(value).strip()
    if not text:
        return None
    # Primero formatos explícitos. Evita que pandas interprete 2026-07-01 como 2026-01-07 con dayfirst=True.
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%d-%m-%Y", "%Y/%m/%d", "%d.%m.%Y"):
        try:
            return datetime.strptime(text, fmt).date().isoformat()
        except ValueError:
            continue
    if pd is not None:
        try:
            if pd.isna(value):
                return None
        except Exception:
            pass
        try:
            # Excel serial y cadenas no estándar.
            dt = pd.to_datetime(value, dayfirst=True, errors="coerce")
            if dt is not None and not pd.isna(dt):
                return dt.date().isoformat()
        except Exception:
            pass
    return None


def detectar_columnas(columnas: Iterable[Any]) -> dict[str, str]:
    """Mapea nombres flexibles del cronograma a campos internos."""
    disponibles = {normalizar_texto(c): str(c) for c in columnas}
    mapping: dict[str, str] = {}
    for target, synonyms in COLUMN_SYNONYMS.items():
        for synonym in synonyms:
            norm = normalizar_texto(synonym)
            if norm in disponibles:
                mapping[target] = disponibles[norm]
                break
        if target in mapping:
            continue
        # Coincidencia parcial controlada para encabezados largos.
        for norm_col, original in disponibles.items():
            if any(normalizar_texto(s) in norm_col for s in synonyms):
                mapping[target] = original
                break
    return mapping


def canonical_modulo(value: Any, tipo_formato: Any = None) -> str:
    raw = normalizar_texto(value) or normalizar_texto(tipo_formato)
    if "rpp" in raw:
        return "RPP"
    if "bienestarina" in raw or "bienesterina" in raw:
        return "Bienestarina"
    if "ram" in raw or "ran" in raw or "asistencia" in raw:
        return "RAM/RAN/Asistencia"
    if "nutric" in raw or "peso" in raw or "talla" in raw:
        return "Nutrición"
    if "talento" in raw:
        return "Talento Humano"
    if "planeacion" in raw or "planeacion pedagogica" in raw or "pedagog" in raw:
        return "Planeación Pedagógica"
    if "coordinador" in raw:
        return "Gestión por Coordinador"
    if "reporte" in raw or "informe" in raw:
        return "Reportes Gerenciales"
    if "cumplimiento" in raw or "icbf" in raw:
        return "Cumplimiento ICBF"
    return str(value or tipo_formato or "General").strip() or "General"


def calcular_estado_color(fecha_limite: Any, estado: Any = "pendiente", hoy: date | None = None) -> tuple[str, str, int | None]:
    """Calcula estado visual y color según fecha límite y estado funcional."""
    estado_norm = normalizar_texto(estado).replace(" ", "_") or "pendiente"
    if estado_norm not in ESTADOS_PERMITIDOS:
        estado_norm = "pendiente"
    if estado_norm in {"entregado", "aprobado"}:
        return estado_norm, "verde", None
    if estado_norm in {"no_aplica", "cerrado"}:
        return estado_norm, "gris", None
    if estado_norm == "rechazado":
        return estado_norm, "rojo", None

    fecha_iso = parse_fecha(fecha_limite)
    if not fecha_iso:
        return "programado", "azul", None
    hoy = hoy or date.today()
    fecha = datetime.strptime(fecha_iso, "%Y-%m-%d").date()
    dias = (fecha - hoy).days
    if dias < 0:
        return "vencido", "rojo", dias
    if dias == 0:
        return "vence_hoy", "rojo", dias
    if dias <= 2:
        return "proximo_a_vencer", "naranja", dias
    if dias <= 5:
        return "proximo_a_vencer", "amarillo", dias
    return "programado", "azul", dias


def fechas_recurrentes(
    fecha_inicial: Any,
    recurrencia: Any = "ninguna",
    hasta: Any = None,
    intervalo: Any = 1,
    *,
    max_instancias: int = 366,
) -> list[str]:
    """Expande una recurrencia finita sin inventar una fecha de terminación."""
    inicial_iso = parse_fecha(fecha_inicial)
    if not inicial_iso:
        raise ValueError("La fecha inicial de la recurrencia es inválida.")
    tipo = normalizar_texto(recurrencia).replace(" ", "_") or "ninguna"
    if tipo not in RECURRENCIAS_PERMITIDAS:
        raise ValueError("Recurrencia no permitida. Usa ninguna, diaria, semanal o mensual.")
    if tipo == "ninguna":
        return [inicial_iso]
    hasta_iso = parse_fecha(hasta)
    if not hasta_iso:
        raise ValueError("recurrencia_hasta es obligatoria para una actividad recurrente.")
    start = datetime.strptime(inicial_iso, "%Y-%m-%d").date()
    end = datetime.strptime(hasta_iso, "%Y-%m-%d").date()
    if end < start:
        raise ValueError("recurrencia_hasta no puede ser anterior a fecha_limite.")
    try:
        step = int(intervalo or 1)
    except (TypeError, ValueError) as exc:
        raise ValueError("recurrencia_intervalo debe ser un número entero.") from exc
    if step < 1 or step > 52:
        raise ValueError("recurrencia_intervalo debe estar entre 1 y 52.")

    dates: list[str] = []
    index = 0
    while len(dates) < max_instancias:
        if tipo == "diaria":
            current = start + timedelta(days=index * step)
        elif tipo == "semanal":
            current = start + timedelta(days=index * step * 7)
        else:
            month_index = (start.month - 1) + index * step
            year = start.year + month_index // 12
            month = month_index % 12 + 1
            current = date(year, month, min(start.day, calendar.monthrange(year, month)[1]))
        if current > end:
            break
        dates.append(current.isoformat())
        index += 1
    if len(dates) >= max_instancias:
        raise ValueError(f"La recurrencia supera el máximo de {max_instancias} instancias.")
    return dates


def clave_unica_entregable(data: dict[str, Any]) -> str:
    parts = [
        parse_fecha(data.get("fecha_limite")) or "",
        normalizar_texto(data.get("titulo")),
        normalizar_texto(data.get("modulo")),
        normalizar_texto(data.get("coordinador")),
        normalizar_texto(data.get("unidad")),
    ]
    base = "|".join(parts)
    return hashlib.sha1(base.encode("utf-8")).hexdigest()




def _dataframe_from_plain_text(texto: str):
    """Convierte texto tabular de Word/PDF/PowerPoint/OCR a DataFrame."""
    if pd is None:
        raise ValueError("Pandas no está disponible para procesar cronogramas.")
    lineas = [linea.strip() for linea in str(texto or "").splitlines() if linea and linea.strip()]
    if not lineas:
        raise ValueError("No se encontró texto legible con estructura de cronograma.")
    muestra = "\n".join(lineas)
    from io import StringIO
    for sep in ["\t", ";", ",", "|"]:
        if sep in muestra:
            try:
                df = pd.read_csv(StringIO(muestra), sep=sep, dtype=str, engine="python")
                if df is not None and not df.empty and len(df.columns) > 1:
                    return df
            except Exception:
                pass
    try:
        df = pd.read_csv(StringIO(muestra), sep=r"\s{2,}", dtype=str, engine="python")
        if df is not None and not df.empty and len(df.columns) > 1:
            return df
    except Exception:
        pass
    # Último recurso: extraer eventos con fecha al inicio o dentro de la línea.
    filas = []
    patron_fecha = re.compile(r"(\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}[/-]\d{1,2}[/-]\d{1,2})")
    for linea in lineas:
        m = patron_fecha.search(linea)
        if not m:
            continue
        fecha = m.group(1)
        actividad = (linea[:m.start()] + " " + linea[m.end():]).strip(" -:|\t")
        if actividad:
            filas.append({"Fecha": fecha, "Actividad": actividad})
    if filas:
        return pd.DataFrame(filas)
    raise ValueError("No se pudo detectar una tabla o fechas dentro del archivo cargado.")


def _dataframe_from_docx(path: str):
    """Lee cronogramas o listas de chequeo Word sin inventar fechas.

    Las listas institucionales suelen traer ``N° / ACTIVIDAD / TH A CARGO /
    ENTREGA`` y filas separadoras de componente. Se conservan todas las filas
    como propuestas; las que no tienen fecha se marcan posteriormente para
    revisión humana en lugar de descartarse.
    """
    try:
        from docx import Document
    except Exception as exc:
        raise ValueError("Para leer Word se requiere python-docx instalado.") from exc
    doc = Document(path)
    frames = []
    for table in doc.tables:
        raw_rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        raw_rows = [row for row in raw_rows if any(str(cell or "").strip() for cell in row)]
        if len(raw_rows) < 2:
            continue
        header = [str(value or "").strip() or f"COLUMNA_{idx + 1}" for idx, value in enumerate(raw_rows[0])]
        if len(set(normalizar_texto(h) for h in header if h)) < 2:
            continue
        mapping = detectar_columnas(header)
        if "titulo" not in mapping:
            continue
        current_component = ""
        records = []
        for values in raw_rows[1:]:
            padded = list(values) + [""] * max(0, len(header) - len(values))
            record = {header[idx]: str(padded[idx] or "").strip() for idx in range(len(header))}
            nonempty = [normalizar_texto(value) for value in padded if str(value or "").strip()]
            # Filas combinadas como "COMPONENTE PEDAGÓGICO" aparecen repetidas
            # en todas las celdas. Se usan como contexto y no como actividad.
            unique_nonempty = set(nonempty)
            if unique_nonempty and len(unique_nonempty) == 1 and "componente" in next(iter(unique_nonempty)):
                current_component = str(next(value for value in padded if str(value or "").strip())).strip()
                continue
            title_column = mapping.get("titulo")
            if not title_column or not str(record.get(title_column) or "").strip():
                continue
            record["COMPONENTE"] = current_component
            records.append(record)
        if records:
            frames.append(pd.DataFrame(records))
    if frames:
        return pd.concat(frames, ignore_index=True, sort=False)
    texto = "\n".join([p.text for p in doc.paragraphs if p.text and p.text.strip()])
    return _dataframe_from_plain_text(texto)


def _dataframe_from_pdf(path: str):
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise ValueError("Para leer PDF se requiere pypdf instalado.") from exc
    reader = PdfReader(path)
    texto = "\n".join(page.extract_text() or "" for page in reader.pages)
    if texto and texto.strip():
        return _dataframe_from_plain_text(texto)
    # Fallback controlado para PDF escaneado. No rompe si falta pdf2image/poppler.
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except Exception as exc:
        raise ValueError(
            "El PDF parece escaneado y no contiene texto extraíble. Para leerlo por OCR se requiere "
            "pdf2image + pytesseract + Poppler/Tesseract instalados. También puedes subir Excel/Word o una imagen legible."
        ) from exc
    textos = []
    for page_img in convert_from_path(path, dpi=220):
        try:
            textos.append(pytesseract.image_to_string(page_img, lang="spa+eng"))
        except Exception:
            textos.append(pytesseract.image_to_string(page_img))
    return _dataframe_from_plain_text("\n".join(textos))


def _dataframe_from_pptx(path: str):
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ValueError("Para leer PowerPoint se requiere python-pptx instalado.") from exc
    prs = Presentation(path)
    textos = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                textos.append(shape.text)
    return _dataframe_from_plain_text("\n".join(textos))


def _dataframe_from_image(path: str):
    try:
        from PIL import Image
        import pytesseract
    except Exception as exc:
        raise ValueError(
            "La imagen fue recibida, pero para extraer cronogramas desde fotografía se requiere OCR "
            "(Pillow + pytesseract + motor Tesseract). Sube Excel/Word/PDF o instala OCR para leer imágenes."
        ) from exc
    imagen = Image.open(path)
    try:
        texto = pytesseract.image_to_string(imagen, lang="spa+eng")
    except Exception:
        # Algunos equipos tienen Tesseract instalado sin paquete de idioma español.
        # Se reintenta con el idioma por defecto para no bloquear el menú.
        texto = pytesseract.image_to_string(imagen)
    return _dataframe_from_plain_text(texto)


def leer_cronograma_flexible(path: str, filename: str = ""):
    """Lee cronogramas desde Excel, CSV, Word, PDF, PowerPoint o imagen.

    Esta función amplía Alpha16 para que el botón Cargar cronograma mensual no
    dependa exclusivamente de Excel. En formatos visuales se intenta extraer texto;
    si no hay OCR, devuelve un error claro sin romper la plataforma.
    """
    if pd is None:
        raise ValueError("Pandas no está disponible para leer cronogramas.")
    name = str(filename or path)
    ext = name.lower().rsplit('.', 1)[-1] if '.' in name else ''
    ext = f'.{ext}' if ext else ''
    if ext in {'.xlsx', '.xls', '.xlsm', '.ods'}:
        hojas = pd.read_excel(path, sheet_name=None, dtype=str)
        mejor = None
        mejor_score = -1
        for _nombre, df in hojas.items():
            if df is None or df.empty:
                continue
            mapping = detectar_columnas(df.columns)
            score = len(mapping) * 100 + len(df)
            if 'fecha_limite' in mapping and 'titulo' in mapping:
                score += 1000
            if score > mejor_score:
                mejor_score = score
                mejor = df
        if mejor is None:
            raise ValueError("El Excel no contiene hojas con datos de cronograma.")
        return mejor
    if ext in {'.csv', '.txt', '.tsv', '.tab', '.dat'}:
        if ext in {'.tsv', '.tab'}:
            return pd.read_csv(path, sep='\t', dtype=str, engine='python')
        try:
            return pd.read_csv(path, sep=None, dtype=str, engine='python')
        except Exception:
            return _dataframe_from_plain_text(open(path, 'r', encoding='utf-8', errors='ignore').read())
    if ext == '.docx':
        return _dataframe_from_docx(path)
    if ext == '.pdf':
        return _dataframe_from_pdf(path)
    if ext == '.pptx':
        return _dataframe_from_pptx(path)
    if ext in {'.png', '.jpg', '.jpeg', '.webp', '.bmp', '.tif', '.tiff'}:
        return _dataframe_from_image(path)
    raise ValueError(f"Extensión {ext or 'desconocida'} no soportada para cronograma.")

def row_to_payload(row: Any, mapping: dict[str, str]) -> dict[str, Any]:
    payload: dict[str, Any] = {}
    for field, col in mapping.items():
        value = row.get(col) if hasattr(row, "get") else None
        if field in {"fecha_limite", "fecha_inicio"}:
            value = parse_fecha(value)
        if value is not None and str(value).strip() != "":
            payload[field] = value
    if not payload.get("titulo"):
        payload["titulo"] = payload.get("descripcion") or payload.get("tipo_formato") or "Entregable operativo"
    payload["modulo"] = canonical_modulo(payload.get("modulo"), payload.get("tipo_formato"))
    payload.setdefault("estado", "pendiente")
    payload.setdefault("prioridad", "Media")
    payload["fecha_limite"] = parse_fecha(payload.get("fecha_limite"))
    payload["fecha_inicio"] = parse_fecha(payload.get("fecha_inicio")) or payload.get("fecha_limite")
    estado, color, dias = calcular_estado_color(payload.get("fecha_limite"), payload.get("estado"))
    payload["estado"] = estado
    payload["color"] = color
    payload["dias_restantes"] = dias
    payload["clave_unica"] = clave_unica_entregable(payload)
    return payload



def construir_preview_cronograma(path: str, filename: str = "") -> dict[str, Any]:
    """Lee un cronograma y devuelve una vista previa editable sin guardar actividades.

    Esta función implementa ALPHA33: primero extrae actividades y fechas desde el
    archivo cargado, normaliza columnas y marca advertencias para que el usuario
    revise antes de publicar en el calendario. No escribe en base de datos.
    """
    df = leer_cronograma_flexible(path, filename)
    mapping = detectar_columnas(df.columns)
    ext = str(filename or path).lower().rsplit('.', 1)[-1] if '.' in str(filename or path) else ''
    requiere_revision = ext in {"pdf", "png", "jpg", "jpeg", "webp", "bmp", "tif", "tiff"}
    actividades: list[dict[str, Any]] = []
    errores: list[dict[str, Any]] = []

    if "titulo" not in mapping:
        return {
            "total_filas": int(len(df)),
            "actividades": [],
            "errores": [{
                "fila": 0,
                "error": "No se detectó una columna de actividad/entregable.",
                "columnas": [str(c) for c in df.columns],
            }],
            "advertencias": ["Revisa el archivo: no fue posible identificar las actividades."],
            "columnas_detectadas": mapping,
            "requiere_revision": True,
        }

    claves_vistas: set[str] = set()
    duplicados = 0
    for idx, row in df.iterrows():
        fila = int(idx) + 2
        try:
            payload = row_to_payload(row, mapping)
            actividad_errores = []
            advertencias = []
            if not payload.get("fecha_limite"):
                actividad_errores.append("Fecha inválida o vacía.")
            if not str(payload.get("titulo") or "").strip():
                actividad_errores.append("Actividad vacía.")
            clave = payload.get("clave_unica") or clave_unica_entregable(payload)
            if clave in claves_vistas:
                duplicados += 1
                advertencias.append("Posible duplicado en el mismo archivo.")
            claves_vistas.add(clave)
            if requiere_revision:
                advertencias.append("Revisar antes de guardar: la lectura proviene de PDF/imagen/OCR.")
            confianza = 100
            if requiere_revision:
                confianza -= 25
            if not payload.get("fecha_limite"):
                confianza -= 35
            if not payload.get("responsable_nombre"):
                confianza -= 10
            if advertencias:
                confianza -= min(15, len(advertencias) * 5)
            actividades.append({
                "id_temp": len(actividades) + 1,
                "fila_original": fila,
                "fecha": payload.get("fecha_limite") or "",
                "fecha_limite": payload.get("fecha_limite") or "",
                "titulo": payload.get("titulo") or "",
                "descripcion": payload.get("descripcion") or "",
                "componente": payload.get("componente") or "",
                "entregables": payload.get("entregables") or payload.get("tipo_formato") or "",
                "numero": payload.get("numero") or "",
                "responsable_nombre": payload.get("responsable_nombre") or "",
                "coordinador": payload.get("coordinador") or "",
                "unidad": payload.get("unidad") or "",
                "modulo": payload.get("modulo") or "General",
                "tipo_formato": payload.get("tipo_formato") or payload.get("modulo") or "General",
                "estado": payload.get("estado") or "programado",
                "prioridad": payload.get("prioridad") or "Media",
                "observaciones": payload.get("observaciones") or payload.get("descripcion") or "",
                "municipio": payload.get("municipio") or "",
                "clave_unica": clave,
                "ok": not actividad_errores,
                "errores": actividad_errores,
                "advertencias": advertencias,
                "confianza": max(0, min(100, confianza)),
                "origen": filename or str(path),
            })
            if actividad_errores:
                errores.append({"fila": fila, "error": "; ".join(actividad_errores)})
        except Exception as exc:
            errores.append({"fila": fila, "error": str(exc)})

    validas = sum(1 for a in actividades if a.get("ok"))
    return {
        "total_filas": int(len(df)),
        "actividades": actividades,
        "validas": validas,
        "invalidas": len(actividades) - validas,
        "duplicados_en_archivo": duplicados,
        "errores": errores[:100],
        "advertencias": (["El archivo requiere revisión manual antes de guardar."] if requiere_revision else []),
        "columnas_detectadas": mapping,
        "requiere_revision": requiere_revision or bool(errores),
    }
